import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import os
import math

# Set up command line argument parsing
parser = argparse.ArgumentParser(
    description="Extract icons and labels from a PowerPoint file and save them as images."
)
parser.add_argument(
    "filename", 
    help="Filename of the PowerPoint (.pptx) file to process"
)
parser.add_argument(
    "-o", "--output_folder", 
    default="out", 
    help="Output folder to save extracted images (default: 'out')"
)
args = parser.parse_args()

# Get the paths from arguments
pptx_path = args.filename
output_folder = args.output_folder

# Ensure the output folder exists
os.makedirs(output_folder, exist_ok=True)

# Load the presentation
prs = Presentation(pptx_path)

def calculate_distance(shape1, shape2):
    """Calculate the Euclidean distance between two shapes based on their top-left coordinates."""
    x1, y1 = shape1.left, shape1.top
    x2, y2 = shape2.left, shape2.top
    return math.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2)

# Iterate over each slide
for slide_number, slide in enumerate(prs.slides):
    # Process each icon (picture) shape
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Save the icon as an image file
            image_stream = io.BytesIO(shape.image.blob)
            image = Image.open(image_stream)

            # Find the closest text label to this icon
            closest_label = None
            min_distance = float('inf')  # Initialize with infinity

            for text_shape in slide.shapes:
                if text_shape.has_text_frame and text_shape.text_frame.text.strip():
                    distance = calculate_distance(shape, text_shape)

                    # Update closest label if this one is nearer
                    if distance < min_distance:
                        min_distance = distance
                        closest_label = text_shape.text_frame.text.strip()

            # Define a filename based on the closest label, or default if not found
            filename = closest_label if closest_label else f"icon_slide{slide_number+1}_{shape.shape_id}"
            sanitized_filename = "".join(c for c in filename if c.isalnum() or c in " _-")
            filepath = os.path.join(output_folder, f"{sanitized_filename}.png")

            # Save the image
            image.save(filepath)
            print(f"Saved {filepath}")

print("Icon extraction complete!")