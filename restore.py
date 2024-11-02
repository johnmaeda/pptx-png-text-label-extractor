import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from PIL import Image

# Set up command line argument parsing
parser = argparse.ArgumentParser(
    description="Generate a PowerPoint slide deck with icons and labels from a specified directory."
)
parser.add_argument(
    "-d", "--directory",
    default="out",
    help="Directory containing the icon images (default: 'out')"
)
parser.add_argument(
    "-o", "--output",
    default="alltheicons.pptx",
    help="Filename for the output PowerPoint file (default: 'alltheicons.pptx')"
)
args = parser.parse_args()

# Assign the directory and output file from arguments
image_folder = args.directory
output_pptx = args.output

# Create a new PowerPoint presentation
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout

# Define maximum width and grid layout settings
max_icon_width = Inches(0.5)  # Maximum icon width
text_box_width = Inches(1.5)  # Fixed width for text box
icons_per_row = 4  # Number of icons per row
top_margin = Inches(1)
left_margin = Inches(1)
col_offset = Inches(1.2)  # Space between icons
vertical_offset = Inches(0.3)  # Offset for stagger effect
row_offset = Inches(0.8)  # Offset for every other row

# Ensure the image folder exists
if not os.path.exists(image_folder):
    print(f"Error: The directory '{image_folder}' does not exist.")
    exit(1)

# Load all image files in the folder
icon_files = [f for f in os.listdir(image_folder) if f.endswith('.png')]

# Add a new slide and set background color to 50% gray (RGB 128, 128, 128)
def add_slide_with_gray_background(prs):
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(128, 128, 128)  # Set to 50% gray
    return slide

# Create the initial slide
slide = add_slide_with_gray_background(prs)

# Initialize positioning
current_row, current_col = 0, 0
for icon_file in icon_files:
    # Calculate horizontal position
    left_position = left_margin + (current_col * (max_icon_width + col_offset))
    
    # Stagger every other icon vertically
    top_position = top_margin + (current_row * (max_icon_width + row_offset))
    if current_col % 2 != 0:
        top_position += vertical_offset  # Offset every other column

    # Add the image to the slide
    img_path = os.path.join(image_folder, icon_file)
    icon = Image.open(img_path)
    aspect_ratio = icon.width / icon.height
    icon_width = max_icon_width
    icon_height = icon_width / aspect_ratio
    pic = slide.shapes.add_picture(img_path, left_position, top_position, width=icon_width, height=icon_height)

    # Calculate centered position for the text box under the icon
    text_left_position = left_position + (icon_width - text_box_width) / 2
    text_top_position = top_position + icon_height + Inches(0.1)  # Position text just below icon

    # Add the text box with word wrap enabled and centered under the icon
    icon_name = os.path.splitext(icon_file)[0]  # Get the filename without the extension
    text_box = slide.shapes.add_textbox(text_left_position, text_top_position, text_box_width, Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True  # Enable word wrap
    text_frame.text = icon_name
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color for text
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center text within the box

    # Update column and row for grid layout
    current_col += 1
    if current_col >= icons_per_row:
        current_col = 0
        current_row += 1
        # If there's no space left on the current slide, add a new slide with gray background
        if top_position + 2 * (max_icon_width + row_offset) > prs.slide_height:
            slide = add_slide_with_gray_background(prs)
            current_row = 0

# Save the presentation
prs.save(output_pptx)
print(f"Presentation saved as {output_pptx}")
